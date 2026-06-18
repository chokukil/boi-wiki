from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from scripts.collect_poc_evidence import build_capture_targets
from scripts.check_poc_delivery_readiness import DEFAULT_ARTIFACT_PPTX, evaluate_e2e_summary, evaluate_final_deck
from scripts.insert_poc_screenshots import load_manifest, missing_screenshots, screenshot_issues


def tiny_png_bytes() -> bytes:
    return (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
        b"\x00\x00\x00\x01\x08\x06\x00\x00\x00\x1f\x15\xc4"
        b"\x89\x00\x00\x00\nIDATx\x9cc\xf8\x0f\x00\x01\x01"
        b"\x01\x00\x18\xdd\x8d\xb0\x00\x00\x00\x00IEND\xaeB`\x82"
    )


def test_capture_manifest_lists_required_executive_screenshots():
    manifest = load_manifest(Path("artifacts/boi-poc/capture-manifest.json"))

    assert manifest["capture_dir"] == "captures/boi-poc"
    assert len(manifest["required"]) == 8
    assert [entry["file"] for entry in manifest["required"][:3]] == [
        "01-boi-wiki-home.png",
        "02-sop-library.png",
        "03-event-type-catalog.png",
    ]
    assert any(entry["id"] == "langflow" and "7860" in entry["url"] for entry in manifest["required"])
    assert any(entry["id"] == "kafka_ui" and "8081" in entry["url"] for entry in manifest["required"])


def test_insert_script_detects_missing_screenshots_before_final_output():
    manifest = load_manifest(Path("artifacts/boi-poc/capture-manifest.json"))

    missing = missing_screenshots(manifest)

    assert missing
    assert missing[0].name == "01-boi-wiki-home.png"


def test_screenshot_check_rejects_tiny_or_invalid_capture_files(tmp_path):
    capture_dir = tmp_path / "captures"
    capture_dir.mkdir()
    (capture_dir / "tiny.png").write_bytes(tiny_png_bytes())
    (capture_dir / "not-png.png").write_text("not a png", encoding="utf-8")
    manifest = {
        "capture_dir": str(capture_dir),
        "deck_input": "unused.pptx",
        "deck_output": "unused-output.pptx",
        "required": [
            {"id": "tiny", "file": "tiny.png", "title": "Tiny", "url": "http://localhost", "purpose": "tiny"},
            {"id": "bad", "file": "not-png.png", "title": "Bad", "url": "http://localhost", "purpose": "bad"},
        ],
    }

    issues = screenshot_issues(manifest)

    assert any(issue["id"] == "tiny" and "too small" in issue["reason"] for issue in issues)
    assert any(issue["id"] == "bad" and "not a valid PNG" in issue["reason"] for issue in issues)


def test_executive_ppt_contains_capture_slot_and_technical_appendix():
    prs = Presentation("artifacts/boi-poc/boi-wiki-poc-executive-brief.pptx")
    all_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

    assert len(prs.slides) == 15
    assert "실제 화면 캡처 삽입 계획" in all_text
    assert "Appendix A. 기술 구성" in all_text
    assert "google/gemma-4-26b-a4b-qat" in all_text


def test_capture_targets_resolve_latest_private_boi_and_langflow_urls():
    evidence = {
        "collected_at": "2026-06-17T00:00:00+00:00",
        "git_commit": "abc123",
        "boi_docs": {
            "items": [
                {
                    "event_type": "corrective_action.requested.v1",
                    "uri": "/private/100001/boi-private-corrective.md",
                }
            ]
        },
        "langflow_smoke": {"flow": {"id": "flow-123"}},
    }

    targets = build_capture_targets(
        evidence,
        Path("artifacts/boi-poc/capture-manifest.json"),
        "http://localhost:8000",
        "http://localhost:7860",
    )
    by_id = {target["id"]: target for target in targets["targets"]}

    assert by_id["private_boi"]["url"] == "http://localhost:8000/docs/private/100001/boi-private-corrective.md"
    assert by_id["langflow"]["url"] == "http://localhost:7860/flow/flow-123"
    assert targets["artifact_deck_output"].endswith(
        "outputs/manual-20260619/presentations/boi-e2e-evidence/output/boi-wiki-e2e-evidence-brief.pptx"
    )
    assert targets["legacy_screenshot_deck_output"] == "artifacts/boi-poc/boi-wiki-poc-executive-brief-with-screenshots.pptx"
    assert "<latest" not in by_id["private_boi"]["url"]
    assert "<latest" not in by_id["langflow"]["url"]


def test_e2e_ppt_build_script_preserves_artifact_tool_boundary():
    script = Path("scripts/build_boi_e2e_ppt.py").read_text(encoding="utf-8")
    status = Path("artifacts/boi-poc/ppt-production-status.md").read_text(encoding="utf-8")

    assert "check_presentation_runtime.mjs" in script
    assert "build_artifact_deck.mjs" in script
    assert "--slide-count" in script
    assert "EXPECTED_SLIDE_COUNT = 8" in script
    assert "python scripts/build_boi_e2e_ppt.py" in status
    assert "artifact-tool runtime preflight" in status


def test_capture_targets_point_to_latest_sso_e2e_trace():
    manifest = Path("artifacts/boi-poc/capture-manifest.json").read_text(encoding="utf-8")
    targets = Path("artifacts/boi-poc/capture-targets.json").read_text(encoding="utf-8")
    docs_manifest = Path("docs/POC_SCREENSHOT_MANIFEST.md").read_text(encoding="utf-8")
    combined = "\n".join([manifest, targets, docs_manifest])

    assert "trace-609660cf137c4946aaa833c891f704b7" in combined
    assert "boi:private:100001:20260619014436:7ff90d" in combined
    assert "422fa3e4-d09b-4d51-b323-e652a13f2792" in combined
    assert "3aba3309-89a8-4171-a153-00db6b16dcba" not in combined
    assert "boi-private-100001-20260618012251-15654c.md" not in combined


def test_capture_target_preflight_checks_boi_langflow_and_kafka_urls():
    script = Path("scripts/check_poc_capture_targets.py").read_text(encoding="utf-8")

    assert "x-service-token" in script
    assert "/api/v1/flows/" in script
    assert "flow id not found in Langflow API" in script
    assert "trace-609660cf137c4946aaa833c891f704b7" in script
    assert "boi:private:100001:20260619014436:7ff90d" in script


def test_capture_docs_name_artifact_tool_deck_as_final_delivery():
    capture_targets = Path("artifacts/boi-poc/capture-targets.md").read_text(encoding="utf-8")
    screenshot_manifest = Path("docs/POC_SCREENSHOT_MANIFEST.md").read_text(encoding="utf-8")
    combined = "\n".join([capture_targets, screenshot_manifest])

    assert "Artifact-tool final deck" in combined
    assert "outputs/manual-20260619/presentations/boi-e2e-evidence/output/boi-wiki-e2e-evidence-brief.pptx" in combined
    assert "Legacy screenshot insertion output" in combined
    assert "- Final deck output: `artifacts/boi-poc/boi-wiki-poc-executive-brief-with-screenshots.pptx`" not in combined


def test_delivery_readiness_summary_requires_full_e2e_evidence():
    summary = {
        "trace_id": "trace-ok",
        "event_count": 24,
        "action_count": 21,
        "generated_doc_count": 4,
        "manual_handoff_count": 5,
        "failed_count": 0,
        "langflow_actions": [
            {"action_key": "langflow.boi.reference_flow", "status": "langflow_invoked"},
            {"action_key": "langflow.equipment.stage_analysis", "status": "langflow_invoked"},
        ],
    }

    ok_report = evaluate_e2e_summary(summary)
    bad_report = evaluate_e2e_summary(
        {
            **summary,
            "failed_count": 1,
            "generated_doc_count": 0,
            "langflow_actions": [{"action_key": "langflow.boi.reference_flow", "status": "failed"}],
        }
    )

    assert ok_report["ok"] is True
    assert ok_report["trace_id"] == "trace-ok"
    assert bad_report["ok"] is False
    assert any("failed action" in blocker for blocker in bad_report["blockers"])
    assert any("generated BoI" in blocker for blocker in bad_report["blockers"])
    assert any("Langflow" in blocker for blocker in bad_report["blockers"])


def test_delivery_readiness_script_composes_existing_capture_and_ppt_gates():
    script = Path("scripts/check_poc_delivery_readiness.py").read_text(encoding="utf-8")

    assert "check_poc_capture_targets.py" in script
    assert "insert_poc_screenshots.py" in script
    assert "build_boi_e2e_ppt.py" in script
    assert "screenshots" in script
    assert "ppt_runtime" in script
    assert "final_deck" in script
    assert "blockers" in script


def test_delivery_readiness_final_deck_targets_artifact_tool_output():
    report = evaluate_final_deck(DEFAULT_ARTIFACT_PPTX)

    assert report["path"].endswith(
        "outputs/manual-20260619/presentations/boi-e2e-evidence/output/boi-wiki-e2e-evidence-brief.pptx"
    )
    if not report["ok"]:
        assert "artifact-tool final deck" in report["blockers"][0]
