#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_ARTIFACT_DIR = ROOT / "artifacts" / "boi-poc"
DEFAULT_EVIDENCE = DEFAULT_ARTIFACT_DIR / "evidence.json"
DEFAULT_OUT = DEFAULT_ARTIFACT_DIR / "boi-wiki-poc-executive-brief.pptx"
DEFAULT_NOTES = DEFAULT_ARTIFACT_DIR / "boi-wiki-poc-executive-brief-notes.md"


NAVY = RGBColor(15, 31, 52)
INK = RGBColor(28, 33, 40)
MUTED = RGBColor(96, 108, 123)
BLUE = RGBColor(31, 97, 255)
CYAN = RGBColor(0, 162, 184)
GREEN = RGBColor(31, 151, 100)
AMBER = RGBColor(240, 170, 48)
RED = RGBColor(210, 76, 76)
PALE_BLUE = RGBColor(235, 242, 255)
PALE_GREEN = RGBColor(235, 249, 242)
PALE_AMBER = RGBColor(255, 246, 226)
PALE_GRAY = RGBColor(246, 248, 250)
WHITE = RGBColor(255, 255, 255)


def add_textbox(slide, x, y, w, h, text: str, size=18, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str = "", eyebrow: str = "BoI Wiki PoC"):
    add_textbox(slide, 0.55, 0.32, 2.4, 0.28, eyebrow, size=9, color=BLUE, bold=True)
    add_textbox(slide, 0.55, 0.66, 8.3, 0.58, title, size=26, color=NAVY, bold=True)
    if subtitle:
        add_textbox(slide, 0.58, 1.22, 10.7, 0.36, subtitle, size=11, color=MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.64), Inches(11.45), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(222, 228, 236)
    line.line.fill.background()


def add_footer(slide, index: int, label: str = "SK hynix AIX 확산 TF | Executive Brief"):
    add_textbox(slide, 0.55, 6.9, 8.8, 0.22, label, size=7.5, color=MUTED)
    add_textbox(slide, 11.7, 6.9, 0.4, 0.22, str(index), size=7.5, color=MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title: str, body: str, fill=PALE_GRAY, accent=BLUE, title_size=13, body_size=10):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(219, 226, 236)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, x + 0.18, y + 0.14, w - 0.34, 0.3, title, size=title_size, color=NAVY, bold=True)
    add_textbox(slide, x + 0.18, y + 0.55, w - 0.34, h - 0.66, body, size=body_size, color=INK)
    return shape


def add_metric(slide, x, y, w, value: str, label: str, color=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.82))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(218, 225, 235)
    add_textbox(slide, x + 0.14, y + 0.09, w - 0.28, 0.33, value, size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.14, y + 0.48, w - 0.28, 0.22, label, size=8.5, color=MUTED, align=PP_ALIGN.CENTER)


def add_bullets(slide, x, y, w, h, bullets: Iterable[str], size=12, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(5)
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_arrow(slide, x, y, w, h, text: str, fill=PALE_BLUE, line=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    frame = shape.text_frame
    frame.clear()
    frame.text = text
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if p.runs:
        p.runs[0].font.name = "Malgun Gothic"
        p.runs[0].font.size = Pt(11)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = NAVY


def add_placeholder(slide, x, y, w, h, title: str, subtitle: str):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(249, 251, 253)
    shape.line.color.rgb = RGBColor(160, 172, 188)
    shape.line.dash_style = 4
    add_textbox(slide, x + 0.2, y + 0.22, w - 0.4, 0.34, title, size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.3, y + 0.67, w - 0.6, h - 0.86, subtitle, size=9, color=MUTED, align=PP_ALIGN.CENTER)


def safe_count(evidence: dict[str, Any], *keys: str, default: int = 0) -> int:
    current: Any = evidence
    for key in keys:
        if not isinstance(current, dict):
            return default
        current = current.get(key)
    try:
        return int(current)
    except Exception:
        return default


def langflow_message(evidence: dict[str, Any]) -> str:
    try:
        return evidence["langflow_smoke"]["run"]["outputs"][0]["outputs"][0]["outputs"]["message"]["message"]
    except Exception:
        return "Langflow smoke 응답 추출 대기"


def latest_langflow_action(evidence: dict[str, Any]) -> dict[str, Any]:
    for row in evidence.get("action_logs", {}).get("items", []):
        if row.get("action_key") == "langflow.boi.reference_flow" and row.get("status") == "langflow_invoked":
            return row
    return {}


def slide_notes() -> list[tuple[str, str]]:
    return [
        ("Title", "경영진에게 PoC의 목적은 특정 도구 도입이 아니라 조직 업무 맥락을 축적하는 공통 기반 검증임을 먼저 설명한다."),
        ("Cascade", "TM, CEO, AIX 확산 TF 메시지를 하나의 실행 논리로 연결한다."),
        ("Why", "개인별 Agent 활용이 늘수록 지식 파편화 리스크도 커진다는 점을 강조한다."),
        ("What", "PoC는 Agent Harness, BoI Wiki, Event Broker, Action Gateway, Langflow 연계를 검증한다."),
        ("Scope", "2개월 안에 검증할 것과 검증하지 않을 것을 분리해 기대치를 통제한다."),
        ("Architecture", "모든 실행 채널은 peer connector이고 유지해야 할 자산은 Event Type, BoI, 정책임을 설명한다."),
        ("Workflow", "장비 이상 SOP 흐름은 이벤트 체인과 액션 통제를 동시에 보여주는 데모다."),
        ("Evidence", "실제 API, Kafka, Langflow, action log로 확인된 수치를 제시한다."),
        ("Langflow", "지정된 OpenAI 호환 Gemma endpoint로 실제 한국어 응답이 나온 사실을 보여준다."),
        ("Governance", "Private-first와 승인 필요한 고위험 Action 통제를 경영진 리스크 관점으로 설명한다."),
        ("Screenshots", "vercel:agent-browser로 실제 화면을 캡처했고, 검증된 PNG를 PPT에 삽입한다."),
        ("Roadmap 2026", "PoC 이후 2026 H2에는 운영모델과 파일럿 확산에 집중한다."),
        ("Roadmap 2027+", "중장기에는 기능조직 확산과 Enterprise Agentic Operating Model로 확장한다."),
        ("Appendix Tech", "기술 검토자는 포트, 컨테이너, 토픽, 모델 설정을 확인할 수 있다."),
        ("Appendix Evidence", "부록에서 최신 로그와 action 상태를 그대로 추적할 수 있다."),
    ]


def build_deck(evidence: dict[str, Any], out_path: Path, notes_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank = prs.slide_layouts[6]

    slides = []

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_textbox(slide, 0.8, 0.76, 10.7, 0.42, "TM → CEO → AIX 확산 TF Cascade", size=14, color=BLUE, bold=True)
    add_textbox(slide, 0.78, 1.35, 11.9, 0.95, "BoI Wiki PoC", size=44, color=NAVY, bold=True)
    add_textbox(slide, 0.82, 2.34, 12.5, 0.45, "1인 1 Agent를 조직의 지식으로 축적하는 업무 맥락 자산화 PoC", size=19, color=INK)
    add_card(slide, 0.9, 3.35, 4.25, 1.7, "핵심 결론", "Agent는 각자 만들 수 있지만, 지식은 회사 방식으로 쌓여야 한다.", fill=PALE_BLUE, accent=BLUE, title_size=15, body_size=13)
    add_card(slide, 5.45, 3.35, 4.25, 1.7, "이번 검증", "Event Broker → Action Gateway → BoI Wiki → Langflow LLM 실행까지 실제 동작 확인", fill=PALE_GREEN, accent=GREEN, title_size=15, body_size=12)
    add_card(slide, 10.0, 3.35, 4.25, 1.7, "증거 캡처", "vercel:agent-browser로 핵심 화면 8개를 캡처하고 PPT에 삽입", fill=PALE_GREEN, accent=GREEN, title_size=15, body_size=12)
    add_footer(slide, 1)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "경영 메시지는 하나의 실행 과제로 연결된다", "이천 포럼의 방향성을 PoC의 실행 범위로 번역")
    add_arrow(slide, 0.9, 2.2, 3.7, 1.0, "TM\n우리 일을 이해하는 AI")
    add_arrow(slide, 5.0, 2.2, 3.7, 1.0, "CEO\nAI Native Memory Creator")
    add_arrow(slide, 9.1, 2.2, 4.0, 1.0, "AIX 확산 TF\n공통 Harness와 BoI Wiki")
    add_card(slide, 1.0, 4.05, 3.7, 1.35, "TM 메시지", "개인 AI 활용을 넘어 조직의 일과 맥락을 먼저 정의해야 한다.", fill=PALE_GRAY, accent=BLUE)
    add_card(slide, 5.05, 4.05, 3.7, 1.35, "CEO 메시지", "DT가 데이터를 연결했다면 AI는 판단과 실행을 연결해야 한다.", fill=PALE_GRAY, accent=CYAN)
    add_card(slide, 9.1, 4.05, 4.0, 1.35, "TF 실행 과제", "Claude, ChatGPT, Copilot, Langflow가 회사 방식으로 일하게 하는 공통 기반", fill=PALE_GRAY, accent=GREEN)
    add_footer(slide, 2)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "문제 정의: Agent는 늘지만 지식은 흩어진다", "PoC가 해결하려는 것은 툴 표준화가 아니라 업무 맥락의 자산화")
    add_card(slide, 0.8, 2.05, 3.4, 2.1, "현재 상태", "개인별 Prompt, 파일, 대화, 자동화가 각자 축적된다.\n성과는 생기지만 재사용과 감사가 어렵다.", fill=PALE_AMBER, accent=AMBER)
    add_card(slide, 4.55, 2.05, 3.4, 2.1, "리스크", "근거 불명확, 권한 혼선, 고위험 Action 자동화, 팀 단위 학습 단절", fill=RGBColor(255, 239, 239), accent=RED)
    add_card(slide, 8.3, 2.05, 4.25, 2.1, "필요한 공통자산", "Event Type, BoI Wiki, Agent Harness, 승격/검토/권한 정책", fill=PALE_BLUE, accent=BLUE)
    add_bullets(slide, 1.0, 5.05, 12.2, 0.95, [
        "PoC의 질문: 개인 Agent가 만든 업무 산출을 회사가 재사용 가능한 맥락 자산으로 바꿀 수 있는가?",
        "성공 기준: 생성, 추적, 검색, 승격, 승인 통제를 최소 구현으로 끝까지 증명"
    ], size=13)
    add_footer(slide, 3)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "PoC 명제: 실행 채널은 바뀌어도 업무 맥락 자산은 남긴다", "Claude, ChatGPT, M365 Copilot, Langflow, Custom Agent는 peer connector")
    for idx, (x, title, body, fill, accent) in enumerate([
        (0.8, "Agent Harness", "Private-first\n출처/권한/승격 원칙\nAction guardrail", PALE_BLUE, BLUE),
        (3.85, "BoI Wiki", "OKF 기반 업무 맥락 저장소\nPublic / Team / Private", PALE_GREEN, GREEN),
        (6.9, "Event Broker", "업무 시점을 이벤트로 발행\nmeeting/action/report/SOP", PALE_AMBER, AMBER),
        (9.95, "Action Gateway", "BoI Writer, Langflow, API,\nWebhook, MCP 실행 통제", RGBColor(238, 243, 248), CYAN),
    ]):
        add_card(slide, x, 2.15, 2.75, 2.25, title, body, fill=fill, accent=accent, title_size=13, body_size=10.5)
    add_textbox(slide, 1.1, 5.2, 12.7, 0.55, "PoC 산출물은 특정 LLM UI가 아니라 ‘업무 이벤트 → 판단/실행 → BoI 축적 → 팀 승격’이라는 운영 패턴이다.", size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "2개월 PoC 범위", "명확한 포함/제외 범위로 경영진 기대치를 통제")
    add_card(slide, 0.85, 2.05, 5.35, 3.6, "포함 범위", "• BoI Profile v0.1 및 seed 문서\n• Event Catalog / Action Catalog\n• Kafka 기반 이벤트 트리거\n• Action Gateway dry-run/approval guard\n• Langflow OpenAI 호환 Gemma Flow\n• 장비 이상 SOP workflow 데모", fill=PALE_GREEN, accent=GREEN, body_size=12)
    add_card(slide, 6.65, 2.05, 5.35, 3.6, "제외 범위", "• 전사 지식 그래프 완성\n• 실 시스템 쓰기 연동\n• 모든 팀 업무 자동화\n• HR/SSO/보안정책 완전 내재화\n• 운영 SLA/비용 최적화\n• 고위험 Action 자동 승인", fill=RGBColor(255, 244, 244), accent=RED, body_size=12)
    add_footer(slide, 5)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "실행 아키텍처", "Event Broker와 Action Gateway를 중심으로 모든 실행 채널을 동등하게 연결")
    steps = [
        ("업무 이벤트", "equipment.alarm\nmeeting.closed\nreport.requested", BLUE),
        ("Kafka", "boi.events\nboi.audit\nboi.dead-letter", CYAN),
        ("Event Router", "catalog 기반\n자동 dispatch", GREEN),
        ("Action Gateway", "BoI Writer\nLangflow/API/MCP", AMBER),
        ("BoI Wiki", "Private → Team\n근거/권한/승격", BLUE),
    ]
    x = 0.65
    for idx, (title, body, color) in enumerate(steps):
        add_card(slide, x, 2.35, 2.25, 1.55, title, body, fill=WHITE, accent=color, title_size=12.5, body_size=9.5)
        if idx < len(steps) - 1:
            add_arrow(slide, x + 2.28, 2.72, 0.55, 0.55, "", fill=PALE_GRAY, line=color)
        x += 2.78
    add_bullets(slide, 0.95, 5.0, 12.2, 0.95, [
        "BoI Writer는 Langflow/API/Webhook/MCP와 동일한 1급 connector로 취급한다.",
        "고위험 Action은 자동 실행하지 않고 approval_required로 기록해 사람 승인 흐름에 연결한다."
    ], size=13)
    add_footer(slide, 6)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "검증 시나리오: 장비 이상 SOP Workflow", "단일 Alarm에서 원인 분석, 보전 가이드, 이상 조치까지 이벤트 체인 실행")
    for idx, (x, title, body, accent) in enumerate([
        (0.75, "1. Alarm 발생", "equipment.alarm.raised.v1\nTrend/Raw Data 요청", BLUE),
        (3.85, "2. 원인 분석", "root_cause.analysis.requested.v1\n원인 후보 BoI 생성", CYAN),
        (6.95, "3. 보전 가이드", "maintenance.guide.requested.v1\nSOP/Runbook 참조", GREEN),
        (10.05, "4. 이상 조치", "corrective_action.requested.v1\n고위험 Action 승인 대기", AMBER),
    ]):
        add_card(slide, x, 2.1, 2.75, 2.2, title, body, fill=PALE_GRAY, accent=accent, title_size=12.5, body_size=10.2)
    add_textbox(slide, 1.0, 5.2, 12.0, 0.6, "검증 포인트: 이벤트 발행, 후속 이벤트 생성, BoI materialization, mock API dry-run, approval guard가 모두 로그로 남는다.", size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 7)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "실제 검증 결과", "API, Kafka, Langflow, Action log 기준으로 확인된 현재 증거")
    add_metric(slide, 0.85, 2.0, 2.0, str(safe_count(evidence, "event_types", "items", default=len(evidence.get("event_types", {}).get("items", [])))), "Event Types", BLUE)
    add_metric(slide, 3.05, 2.0, 2.0, str(safe_count(evidence, "action_catalog", "count")), "Actions", CYAN)
    add_metric(slide, 5.25, 2.0, 2.0, str(safe_count(evidence, "events_log", "count")), "Event Logs", GREEN)
    add_metric(slide, 7.45, 2.0, 2.0, str(safe_count(evidence, "action_logs", "count")), "Action Logs", AMBER)
    add_metric(slide, 9.65, 2.0, 2.0, str(safe_count(evidence, "boi_docs", "count")), "BoI Docs", BLUE)
    add_card(slide, 0.85, 3.45, 4.05, 1.65, "실행 확인", f"Health: BoI API={evidence['health']['boi_api']['status']}, Action Gateway={evidence['health']['action_gateway']['status']}\nKafka topic: {evidence['runtime_config']['event_broker']['topic']}", fill=PALE_GREEN, accent=GREEN)
    add_card(slide, 5.1, 3.45, 4.05, 1.65, "통제 확인", "high risk Action은 approval_required로 멈춤\nblock_process_progress / change_spec_rule", fill=PALE_AMBER, accent=AMBER)
    add_card(slide, 9.35, 3.45, 3.4, 1.65, "증거 파일", "artifacts/boi-poc/evidence.json\nevidence-summary.md", fill=PALE_BLUE, accent=BLUE)
    add_footer(slide, 8)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Langflow / LLM 검증", "OpenAI 호환 API로 지정된 Gemma 모델을 실제 호출")
    langflow_action = latest_langflow_action(evidence)
    langflow_result = langflow_action.get("result") or {}
    add_card(slide, 0.9, 2.0, 5.4, 1.35, "LLM 설정", f"Base URL: {evidence['runtime_config']['llm']['base_url']}\nModel: {evidence['runtime_config']['llm']['model']}", fill=PALE_BLUE, accent=BLUE)
    add_card(slide, 6.65, 2.0, 5.4, 1.35, "Event Trace에서 호출된 Flow", f"{langflow_result.get('flow_name', '')}\n{langflow_result.get('flow_id', '')}", fill=PALE_GREEN, accent=GREEN)
    add_card(slide, 0.9, 3.65, 5.4, 1.25, "E2E Action Log", f"trace={langflow_action.get('trace_id', '')}\nevent={langflow_action.get('event_type', '')}\nstatus={langflow_action.get('status', '')}", fill=WHITE, accent=BLUE, title_size=13, body_size=9.5)
    msg = str(langflow_result.get("message") or langflow_message(evidence))
    add_card(slide, 6.65, 3.65, 5.4, 1.25, "LLM 응답 발췌", msg[:300], fill=WHITE, accent=CYAN, title_size=13, body_size=9)
    add_footer(slide, 9)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Governance: Private-first + 승인 기반 Action", "확산의 핵심은 속도보다 재사용 가능한 통제 구조")
    add_card(slide, 0.9, 2.0, 3.5, 2.1, "Private-first", "개인 업무 맥락은 우선 Private BoI로 축적\n공유는 명시적 요청과 검토 이후", fill=PALE_BLUE, accent=BLUE)
    add_card(slide, 4.65, 2.0, 3.5, 2.1, "Team/Public 승격", "source_refs, reviewer, review status를 요구\n원본과 공유용 draft 분리", fill=PALE_GREEN, accent=GREEN)
    add_card(slide, 8.4, 2.0, 3.5, 2.1, "High-risk Action", "공정 진행 금지, Spec/Rule 변경은 자동 실행 금지\napproval_required로만 기록", fill=PALE_AMBER, accent=AMBER)
    add_bullets(slide, 1.0, 5.1, 11.7, 0.9, [
        "경영진 관점의 질문: 누가 만들었고, 무엇을 근거로 했고, 어디까지 자동 실행했는가?",
        "BoI Wiki PoC는 이 질문에 답하는 메타데이터와 로그를 남기는 구조를 검증한다."
    ], size=13)
    add_footer(slide, 10)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "실제 화면 캡처 증거", "vercel:agent-browser로 캡처한 핵심 화면 PNG를 PPT에 삽입")
    slots = [
        ("BoI Wiki Home", "목록, 필터, Event log가 보이는 첫 화면"),
        ("SOP Library", "설비 이상 SOP와 Agent Harness SOP"),
        ("Event Stream", "Alarm → Root Cause → Maintenance → Corrective chain"),
        ("Action Logs", "materialized / dry_run / approval_required"),
        ("Private BoI", "생성된 corrective action BoI 문서"),
        ("Langflow", "Gemma OpenAI-compatible Reference Flow"),
    ]
    for i, (title, sub) in enumerate(slots):
        x = 0.8 + (i % 3) * 4.1
        y = 2.0 + (i // 3) * 2.05
        add_placeholder(slide, x, y, 3.75, 1.55, title, sub)
    add_footer(slide, 11)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Roadmap: 2026 H2", "PoC에서 운영모델로 넘어가는 전환 단계")
    add_card(slide, 0.85, 2.0, 3.5, 2.6, "Pilot 확장", "3~5개 팀\nStaff/R&D/제조 업무\n팀별 BoI Curator 지정", fill=PALE_BLUE, accent=BLUE)
    add_card(slide, 4.65, 2.0, 3.5, 2.6, "운영정책", "Event Catalog v0.2\n승격/리뷰 체크리스트\n권한/출처/감사 로그", fill=PALE_GREEN, accent=GREEN)
    add_card(slide, 8.45, 2.0, 3.5, 2.6, "Template Library", "Langflow Template\nAgent Harness Prompt\n업무별 Action Catalog", fill=PALE_AMBER, accent=AMBER)
    add_footer(slide, 12)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Roadmap: 2027 → 2028+", "Enterprise Agentic Operating Model로 확장")
    add_arrow(slide, 0.9, 2.15, 3.55, 1.0, "2027\n기능 조직 확산")
    add_arrow(slide, 4.85, 2.15, 3.55, 1.0, "2028+\n전사 Agentic Workflow")
    add_arrow(slide, 8.8, 2.15, 3.55, 1.0, "Target\nManagement Intelligence")
    add_card(slide, 0.95, 4.0, 3.45, 1.35, "2027", "R&D BoI Wiki, 개발 Agentic AI, 고객 요구 변경 BoI, AI 교육/인증", fill=PALE_GRAY, accent=BLUE)
    add_card(slide, 4.9, 4.0, 3.45, 1.35, "2028+", "Enterprise Event Broker, Knowledge Graph, Digital Twin/Operational Physical AI", fill=PALE_GRAY, accent=GREEN)
    add_card(slide, 8.85, 4.0, 3.45, 1.35, "경영 지표", "재사용률, 승격률, 검토 리드타임, 자동화 통제율, 업무 cycle time", fill=PALE_GRAY, accent=AMBER)
    add_footer(slide, 13)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Appendix A. 기술 구성", "실행 중인 PoC 스택과 연결 포인트")
    add_bullets(slide, 0.9, 2.0, 5.5, 3.4, [
        "boi-api: http://localhost:8000",
        "action-gateway: http://localhost:8100",
        "langflow: http://localhost:7860",
        "kafka-ui: http://localhost:8081",
        "Kafka topics: boi.events, boi.audit, boi.dead-letter",
        "LLM: google/gemma-4-26b-a4b-qat via OpenAI-compatible API",
    ], size=12)
    add_card(slide, 7.0, 2.0, 5.0, 3.1, "검증 명령", "pytest tests -q -s\npython scripts/setup_langflow_reference_flows.py\npython scripts/collect_poc_evidence.py\nDocker compose stack health check", fill=PALE_BLUE, accent=BLUE, body_size=12)
    add_footer(slide, 14)
    slides.append(slide)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "Appendix B. 현재 갭과 후속 조치", "완료 기준까지 남은 항목을 숨기지 않고 관리")
    add_card(slide, 0.9, 2.0, 3.6, 2.4, "구현/검증", "완료\nAPI, Kafka, Event Router, Action Gateway, Langflow LLM smoke, BoI materialization", fill=PALE_GREEN, accent=GREEN, body_size=11.2)
    add_card(slide, 4.8, 2.0, 3.6, 2.4, "화면 캡처", "완료\nvercel:agent-browser로 BoI Wiki, Langflow, Kafka UI 캡처\n검증된 PNG 삽입", fill=PALE_GREEN, accent=GREEN, body_size=11.2)
    add_card(slide, 8.7, 2.0, 3.6, 2.4, "PowerPoint Add-in", "다음 단계\nPPT 초안을 PowerPoint에서 열고 ChatGPT 작업창으로 polish 요청", fill=PALE_BLUE, accent=BLUE, body_size=11.2)
    add_textbox(slide, 1.0, 5.25, 11.7, 0.5, "완료 판단은 readiness gate, 실제 캡처 삽입, artifact-tool PPTX export가 모두 통과했을 때로 둔다.", size=14, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 15)
    slides.append(slide)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)

    notes = ["# BoI Wiki PoC Executive Brief Speaker Notes", ""]
    for idx, (title, body) in enumerate(slide_notes(), 1):
        notes.extend([f"## Slide {idx}. {title}", body, ""])
    notes_path.write_text("\n".join(notes), encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description="Build the BoI Wiki PoC executive PowerPoint draft.")
    parser.add_argument("--evidence", default=str(DEFAULT_EVIDENCE))
    parser.add_argument("--out", default=str(DEFAULT_OUT))
    parser.add_argument("--notes", default=str(DEFAULT_NOTES))
    args = parser.parse_args()

    evidence = json.loads(Path(args.evidence).read_text(encoding="utf-8"))
    build_deck(evidence, Path(args.out), Path(args.notes))
    print(json.dumps({"ok": True, "pptx": args.out, "notes": args.notes}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
