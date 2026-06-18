#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import struct
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_MANIFEST = ROOT / "artifacts" / "boi-poc" / "capture-manifest.json"
MIN_CAPTURE_WIDTH = 800
MIN_CAPTURE_HEIGHT = 600


def resolve_project_path(value: str | Path) -> Path:
    path = Path(value)
    return path if path.is_absolute() else ROOT / path


def load_manifest(path: Path) -> dict[str, Any]:
    manifest = json.loads(path.read_text(encoding="utf-8"))
    required = manifest.get("required")
    if not isinstance(required, list) or not required:
        raise ValueError("capture manifest must include non-empty required list")
    for entry in required:
        for field in ["id", "file", "title", "url", "purpose"]:
            if not entry.get(field):
                raise ValueError(f"capture entry is missing {field}: {entry}")
    return manifest


def missing_screenshots(manifest: dict[str, Any]) -> list[Path]:
    capture_dir = resolve_project_path(manifest["capture_dir"])
    missing: list[Path] = []
    for entry in manifest["required"]:
        path = capture_dir / entry["file"]
        if not path.exists():
            missing.append(path)
    return missing


def png_dimensions(path: Path) -> tuple[int, int]:
    with path.open("rb") as handle:
        header = handle.read(24)
    if len(header) < 24 or header[:8] != b"\x89PNG\r\n\x1a\n" or header[12:16] != b"IHDR":
        raise ValueError("not a valid PNG file")
    width, height = struct.unpack(">II", header[16:24])
    return int(width), int(height)


def screenshot_issues(
    manifest: dict[str, Any],
    *,
    min_width: int = MIN_CAPTURE_WIDTH,
    min_height: int = MIN_CAPTURE_HEIGHT,
) -> list[dict[str, Any]]:
    capture_dir = resolve_project_path(manifest["capture_dir"])
    issues: list[dict[str, Any]] = []
    for entry in manifest["required"]:
        path = capture_dir / entry["file"]
        if not path.exists():
            issues.append({"id": entry.get("id"), "file": str(path), "reason": "missing"})
            continue
        try:
            width, height = png_dimensions(path)
        except ValueError as exc:
            issues.append({"id": entry.get("id"), "file": str(path), "reason": str(exc)})
            continue
        if width < min_width or height < min_height:
            issues.append(
                {
                    "id": entry.get("id"),
                    "file": str(path),
                    "reason": f"too small: {width}x{height}; expected at least {min_width}x{min_height}",
                    "width": width,
                    "height": height,
                }
            )
    return issues


def add_textbox(slide, x, y, w, h, text: str, size=12, color=RGBColor(32, 38, 46), bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def cover_slide(slide, prs: Presentation) -> None:
    cover = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height,
    )
    cover.fill.solid()
    cover.fill.fore_color.rgb = RGBColor(255, 255, 255)
    cover.line.fill.background()


def add_title(slide, title: str, subtitle: str) -> None:
    add_textbox(slide, 0.55, 0.35, 2.5, 0.28, "BoI Wiki PoC", size=9, color=RGBColor(31, 97, 255), bold=True)
    add_textbox(slide, 0.55, 0.72, 9.6, 0.55, title, size=25, color=RGBColor(15, 31, 52), bold=True)
    add_textbox(slide, 0.58, 1.24, 11.0, 0.32, subtitle, size=10.5, color=RGBColor(96, 108, 123))
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.62), Inches(11.45), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(222, 228, 236)
    line.line.fill.background()


def fit_image(slide, image_path: Path, x: float, y: float, w: float, h: float) -> None:
    # python-pptx preserves image aspect ratio when either width or height is omitted.
    # Use a clipped visual frame by centering the image on a white card.
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(248, 250, 252)
    card.line.color.rgb = RGBColor(205, 215, 228)

    pic = slide.shapes.add_picture(str(image_path), Inches(x + 0.04), Inches(y + 0.04), width=Inches(w - 0.08))
    if pic.height > Inches(h - 0.08):
        slide.shapes._spTree.remove(pic._element)
        pic = slide.shapes.add_picture(str(image_path), Inches(x + 0.04), Inches(y + 0.04), height=Inches(h - 0.08))
    pic.left = Inches(x + (w - pic.width / 914400) / 2)
    pic.top = Inches(y + (h - pic.height / 914400) / 2)


def update_screenshot_slide(prs: Presentation, manifest: dict[str, Any], slide_number: int) -> None:
    capture_dir = resolve_project_path(manifest["capture_dir"])
    slide = prs.slides[slide_number - 1]
    cover_slide(slide, prs)
    add_title(slide, "실제 화면 캡처", "검증 화면을 PPT에 직접 삽입한 최종 캡처 슬라이드")

    positions = [
        (0.65, 1.95, 3.85, 1.65),
        (4.75, 1.95, 3.85, 1.65),
        (8.85, 1.95, 3.85, 1.65),
        (0.65, 4.15, 3.85, 1.65),
        (4.75, 4.15, 3.85, 1.65),
        (8.85, 4.15, 3.85, 1.65),
    ]
    for entry, (x, y, w, h) in zip(manifest["required"][:6], positions):
        add_textbox(slide, x, y - 0.28, w, 0.2, entry["title"], size=8.5, color=RGBColor(15, 31, 52), bold=True)
        fit_image(slide, capture_dir / entry["file"], x, y, w, h)

    add_textbox(
        slide,
        0.65,
        6.55,
        12.0,
        0.24,
        "Langflow와 Kafka UI 캡처는 Appendix screenshot slide에 별도 삽입된다.",
        size=8,
        color=RGBColor(96, 108, 123),
        align=PP_ALIGN.CENTER,
    )


def add_appendix_slide(prs: Presentation, manifest: dict[str, Any]) -> None:
    capture_dir = resolve_project_path(manifest["capture_dir"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Appendix C. Platform Screens", "Langflow와 Kafka UI 실행 상태")
    for entry, position in zip(manifest["required"][6:], [(0.85, 2.05, 5.35, 3.35), (6.65, 2.05, 5.35, 3.35)]):
        x, y, w, h = position
        add_textbox(slide, x, y - 0.32, w, 0.24, entry["title"], size=11, color=RGBColor(15, 31, 52), bold=True)
        fit_image(slide, capture_dir / entry["file"], x, y, w, h)
        add_textbox(slide, x, y + h + 0.08, w, 0.32, entry["purpose"], size=8.5, color=RGBColor(96, 108, 123))
    add_textbox(slide, 0.55, 6.9, 8.8, 0.22, "SK hynix AIX 확산 TF | Executive Brief", size=7.5, color=RGBColor(96, 108, 123))
    add_textbox(slide, 11.7, 6.9, 0.4, 0.22, str(len(prs.slides)), size=7.5, color=RGBColor(96, 108, 123), align=PP_ALIGN.RIGHT)


def insert_screenshots(manifest: dict[str, Any], output: Path | None = None, allow_missing: bool = False) -> Path:
    issues = screenshot_issues(manifest)
    if issues and not allow_missing:
        formatted = "\n".join(f"{item['file']}: {item['reason']}" for item in issues)
        raise FileNotFoundError(f"Required screenshots are not ready:\n{formatted}")

    deck_input = resolve_project_path(manifest["deck_input"])
    deck_output = output or resolve_project_path(manifest["deck_output"])
    prs = Presentation(str(deck_input))
    if len(prs.slides) < 11:
        raise ValueError("expected the executive deck to have at least 11 slides")

    if not issues:
        update_screenshot_slide(prs, manifest, slide_number=11)
        add_appendix_slide(prs, manifest)
    deck_output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(deck_output)
    return deck_output


def main() -> None:
    parser = argparse.ArgumentParser(description="Insert captured PoC screenshots into the executive PPT.")
    parser.add_argument("--manifest", default=str(DEFAULT_MANIFEST))
    parser.add_argument("--out", default="")
    parser.add_argument("--allow-missing", action="store_true")
    parser.add_argument("--check", action="store_true", help="Only validate screenshot availability.")
    args = parser.parse_args()

    manifest = load_manifest(resolve_project_path(args.manifest))
    issues = screenshot_issues(manifest)
    if args.check:
        missing = [issue["file"] for issue in issues if issue["reason"] == "missing"]
        invalid = [issue for issue in issues if issue["reason"] != "missing"]
        print(
            json.dumps(
                {"ok": not issues, "missing": missing, "invalid": invalid, "issues": issues},
                ensure_ascii=False,
                indent=2,
            )
        )
        raise SystemExit(1 if issues else 0)

    output = resolve_project_path(args.out) if args.out else None
    deck_output = insert_screenshots(manifest, output=output, allow_missing=args.allow_missing)
    print(json.dumps({"ok": True, "pptx": str(deck_output), "issues": issues}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
